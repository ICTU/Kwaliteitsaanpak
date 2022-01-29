"""Kwaliteitsaanpak PPTX-presentation builder."""

import pathlib
import shutil

from lxml import etree
from pptx import Presentation
from pptx.util import Inches, Pt

import xmltags
from custom_types import TreeBuilderAttributes
from .builder import Builder


class PptxBuilder(Builder):
    """Kwaliteitsaanpak presentation builder."""

    # Slide layouts. These are specific for the reference file
    TITLE_SLIDE = 0
    BULLET_SLIDE = 1
    CONTENT_SLIDE = 4
    CHAPTER_SLIDE = 16

    def __init__(self, filename: pathlib.Path, pptx_reference_filename: pathlib.Path) -> None:
        super().__init__(filename)
        filename.unlink(missing_ok=True)
        shutil.copy(pptx_reference_filename, filename)
        self.presentation = Presentation(filename)
        # Make the title placeholder on the measure slide wider
        measure_master_slide = self.presentation.slide_master.slide_layouts[self.CONTENT_SLIDE]
        measure_master_slide.placeholders[0].width = Inches(10)
        self.current_slide = None
        self.chapter_heading = ""

    def start_element(self, tag: str, attributes: TreeBuilderAttributes) -> None:
        super().start_element(tag, attributes)
        if tag == xmltags.SLIDE:
            self.add_slide(self.BULLET_SLIDE, self.chapter_heading)

    def text(self, tag: str, text: str, attributes: TreeBuilderAttributes) -> None:
        if tag == xmltags.TITLE:
            self.add_slide(self.TITLE_SLIDE, text)
            self.current_slide.shapes.title.text_frame.paragraphs[0].font.size = Pt(60)  # type: ignore
        elif tag == xmltags.PARAGRAPH and (self.in_element(xmltags.FRONTPAGE) or self.in_element(xmltags.SLIDE)):
            self.current_slide.shapes[1].text = text  # type: ignore
            self.remove_bullet(paragraph_index=0)
        elif self.in_element(xmltags.MEASURE) and not self.in_appendix():
            if self.chapter_heading:
                self.add_slide(self.CHAPTER_SLIDE, self.chapter_heading)
                self.chapter_heading = ""
            if tag == xmltags.BOLD:
                self.add_slide(self.CONTENT_SLIDE, text)
                self.current_slide.shapes.title.text_frame.paragraphs[0].font.size = Pt(24)  # type: ignore
            else:
                self.add_text_box()
                self.current_slide.shapes[1].text = text  # type: ignore
        elif (
            tag == xmltags.HEADING
            and self.in_element(xmltags.SECTION, dict(level="1"))
            and not self.in_element(xmltags.SECTION, dict(level="2"))
        ):
            self.chapter_heading = text
        elif tag == xmltags.LIST_ITEM and self.in_element(xmltags.SLIDE):
            text_frame = self.current_slide.shapes[1].text_frame  # type: ignore
            if self.current_slide.shapes[1].text.strip():  # type: ignore
                paragraph = text_frame.add_paragraph()
            else:
                paragraph = text_frame.paragraphs[0]
            paragraph.level = 0
            paragraph.text = text

    def add_slide(self, slide_layout_index: int, title: str) -> None:
        """Add a new slide with the given title to the presentation."""
        slide_layout = self.presentation.slide_layouts[slide_layout_index]
        self.current_slide = self.presentation.slides.add_slide(slide_layout)
        self.current_slide.shapes.title.text = title  # type: ignore

    def add_text_box(self):
        """Add a text box to the current slide."""
        text_box = self.current_slide.shapes.add_textbox(Inches(0.7), Inches(1.6), Inches(12), Inches(6))
        text_box.text_frame.word_wrap = True

    def remove_bullet(self, paragraph_index: int):
        """Remove bullets from the paragraph."""
        no_bullet = etree.Element("{http://schemas.openxmlformats.org/drawingml/2006/main}buNone")
        # pylint: disable=protected-access
        self.current_slide.shapes[1].text_frame.paragraphs[paragraph_index]._pPr.insert(0, no_bullet)  # type: ignore

    def in_appendix(self) -> bool:
        """Return whether the current section is an appendix."""
        return self.in_element(xmltags.SECTION, {xmltags.SECTION_IS_APPENDIX: "y"})

    def end_document(self) -> None:
        """Override to save the presentation."""
        self.presentation.save(self.filename)
